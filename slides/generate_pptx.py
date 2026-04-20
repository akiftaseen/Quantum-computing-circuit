import os
from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR

# --- BRANDING & COLOR PALETTE ---
COLOR_PRIMARY = RGBColor(15, 76, 129)    # #0F4C81 (Deep Blue)
COLOR_TEXT = RGBColor(15, 23, 42)       # #0F172A (Near Black)
COLOR_TEXT_SOFT = RGBColor(51, 65, 85)  # #334155 (Gray-Blue)
COLOR_CHIP_BG = RGBColor(232, 242, 250) # #E8F2FA (Light Blue)
COLOR_BORDER = RGBColor(215, 224, 234)  # #D7E0EA (Soft Border)

# Font names (ensure these are installed on your Mac)
FONT_SANS = 'Plus Jakarta Sans'
FONT_MONO = 'JetBrains Mono'

def add_styled_chip(slide, text):
    """Creates the centered top chip with the dot."""
    width, height = Inches(2.0), Inches(0.4)
    left, top = (prs.slide_width - width) / 2, Inches(0.4)
    
    # Background Box
    chip = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    chip.fill.solid()
    chip.fill.fore_color.rgb = COLOR_CHIP_BG
    chip.line.color.rgb = COLOR_PRIMARY
    chip.line.width = Pt(0.75)
    
    # Small Dot inside chip
    dot_size = Inches(0.08)
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.2), top + (height - dot_size) / 2, dot_size, dot_size)
    dot.fill.solid()
    dot.fill.fore_color.rgb = COLOR_PRIMARY
    dot.line.fill.background()
    
    # Text
    tf = chip.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    p.font.name = FONT_MONO
    p.alignment = PP_ALIGN.CENTER

def add_page_index(slide, index_text):
    """Top right index box."""
    width, height = Inches(1.0), Inches(0.4)
    left, top = prs.slide_width - Inches(1.5), Inches(0.4)
    
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = COLOR_CHIP_BG
    box.line.color.rgb = COLOR_BORDER
    
    tf = box.text_frame
    tf.text = index_text
    p = tf.paragraphs[0]
    p.font.size = Pt(12)
    p.font.name = FONT_MONO
    p.font.color.rgb = COLOR_PRIMARY
    p.alignment = PP_ALIGN.CENTER

def generate_presentation():
    global prs
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5) # 16:9

    with open('slides-pptx.html', 'r', encoding='utf-8') as f:
        soup = BeautifulSoup(f, 'html.parser')

    for slide_html in soup.find_all('section', class_='slide'):
        slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank
        
        # 1. Background Logic (Simplified)
        # To get the orbits, you'd need the PNGs of the orbits. 
        # For now, we keep it clean white as per the main screenshot area.

        # 2. Header Elements
        chip_div = slide_html.find('div', class_='section-chip')
        if chip_div:
            add_styled_chip(slide, chip_div.get_text(strip=True))
            
        index_div = slide_html.find('div', class_='slide-index')
        if index_div:
            add_page_index(slide, index_div.get_text(strip=True))

        # 3. Main Title (Centered)
        title_text = slide_html.find('h2', class_='slide-title').get_text(strip=True)
        title_box = slide.shapes.add_textbox(0, Inches(1.0), prs.slide_width, Inches(1.2))
        tf_title = title_box.text_frame
        p_title = tf_title.paragraphs[0]
        p_title.text = title_text
        p_title.font.size = Pt(54)
        p_title.font.bold = True
        p_title.font.name = FONT_SANS
        p_title.font.color.rgb = COLOR_TEXT
        p_title.alignment = PP_ALIGN.CENTER

        # 4. Content Area (Split 45/55)
        left_content = slide_html.find('div', class_='media-left')
        right_content = slide_html.find('div', class_='media-right')

        # Bullets on the Left
        if left_content:
            tx_bullets = slide.shapes.add_textbox(Inches(0.6), Inches(2.8), Inches(5.5), Inches(4.0))
            tf_bul = tx_bullets.text_frame
            tf_bul.word_wrap = True
            
            # Find all bullets or cards
            items = left_content.find_all(['li', 'h3']) 
            for item in items:
                p = tf_bul.add_paragraph()
                p.text = item.get_text(strip=True)
                p.font.size = Pt(26)
                p.font.name = FONT_SANS
                p.font.color.rgb = COLOR_TEXT_SOFT
                p.font.bold = (item.name == 'h3')
                # Custom Blue Bullet
                p.level = 0
                p.space_before = Pt(12)
                # Note: bullet color is hard to change per-paragraph in python-pptx, 
                # but the text color is correctly set to your theme.

        # Image on the Right
        if right_content:
            img_tag = right_content.find('img')
            if img_tag:
                img_path = img_tag['src'].replace('./', '').replace('%20', ' ')
                if os.path.exists(img_path):
                    # Position the image
                    pic = slide.shapes.add_picture(img_path, Inches(6.5), Inches(2.8), width=Inches(6.2))
                    
                    # Add Caption below image
                    caption = right_content.find('div', class_='photo-caption')
                    if caption:
                        cap_box = slide.shapes.add_textbox(Inches(6.5), Inches(6.8), Inches(6.2), Inches(0.4))
                        tf_cap = cap_box.text_frame
                        p_cap = tf_cap.paragraphs[0]
                        p_cap.text = caption.get_text(strip=True).upper()
                        p_cap.font.size = Pt(11)
                        p_cap.font.name = FONT_MONO
                        p_cap.font.color.rgb = COLOR_TEXT_SOFT
                        p_cap.alignment = PP_ALIGN.CENTER

    prs.save('Quantum_Simulator_Final.pptx')
    print("Presentation Generated: Quantum_Simulator_Final.pptx")

if __name__ == "__main__":
    generate_presentation()